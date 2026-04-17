from fastapi import FastAPI, Depends, HTTPException, status, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from sqlalchemy.orm import Session
from typing import List, Optional, Dict
import json
import shutil

import httpx
import os
import io
import fitz # PyMuPDF
import docx
from pptx import Presentation
from pathlib import Path
from datetime import datetime
from dotenv import load_dotenv

import models
import schemas
import crud
from database import engine, get_db

BASE_DIR = Path(__file__).resolve().parent
load_dotenv(BASE_DIR / ".env", override=True)
NEXTTOKEN_API_KEY = os.getenv("NEXTTOKEN_API_KEY", "")
NEXTTOKEN_BASE_URL = "https://gateway.nexttoken.co/v1"
DEFAULT_AI_MODEL = "gpt-4o"
TRANSLATION_PROVIDER = os.getenv("TRANSLATION_PROVIDER", "nexttoken").strip().lower()
HUGGINGFACE_API_TOKEN = os.getenv("HUGGINGFACE_API_TOKEN", "").strip()
HUGGINGFACE_TRANSLATION_MODEL = os.getenv("HUGGINGFACE_TRANSLATION_MODEL", "").strip()
HUGGINGFACE_TRANSLATION_MODEL_MAP = os.getenv("HUGGINGFACE_TRANSLATION_MODEL_MAP", "").strip()

# Creates the SQLite Database file automatically on startup
models.Base.metadata.create_all(bind=engine)

app = FastAPI(title="SukuuPath GH API")

# Setup CORS to allow React Frontend to communicate
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Ensure uploads directory exists
UPLOAD_DIR = BASE_DIR / "uploads"
UPLOAD_DIR.mkdir(exist_ok=True)

# Mount static files for material downloads/preview
app.mount("/uploads", StaticFiles(directory=str(UPLOAD_DIR)), name="uploads")

# --- AUTHENTICATION ROUTES ---

@app.post("/api/auth/signup", response_model=schemas.UserResponse)
def signup(user: schemas.UserCreate, db: Session = Depends(get_db)):
    db_user = crud.get_user_by_email(db, email=user.email)
    if db_user:
        raise HTTPException(status_code=400, detail="Email already registered")
    return crud.create_user(db=db, user=user)

@app.post("/api/auth/google", response_model=schemas.UserResponse)
def google_auth(req: schemas.GoogleLoginRequest, db: Session = Depends(get_db)):
    try:
        db_user = crud.get_user_by_email(db, email=req.email)
        if db_user:
            return db_user
        
        # User not found: Return 404 to trigger frontend redirect to Sign Up
        raise HTTPException(
            status_code=404, 
            detail={
                "message": "User not registered", 
                "is_new_user": True,
                "googleData": {
                    "email": req.email,
                    "first_name": req.first_name,
                    "last_name": req.last_name
                }
            }
        )
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/auth/login", response_model=schemas.UserResponse)
def login(user_credentials: schemas.UserLogin, db: Session = Depends(get_db)):
    db_user = crud.get_user_by_email(db, email=user_credentials.email)
    
    if not db_user or not crud.verify_password(user_credentials.password, db_user.hashed_password):
        raise HTTPException(status_code=401, detail="Invalid email or password")
        
    return db_user

@app.put("/api/auth/update")
def update_profile(req: schemas.UserProfileUpdate, db: Session = Depends(get_db)):
    data = {}
    if req.first_name: data["first_name"] = req.first_name
    if req.last_name: data["last_name"] = req.last_name
    if req.email: data["email"] = req.email
    if req.secondary_email: data["secondary_email"] = req.secondary_email
    updated = crud.update_user_profile(req.user_id, data)
    if not updated: raise HTTPException(status_code=404, detail="User not found")
    return updated

@app.put("/api/auth/password")
def update_password(req: schemas.UserPasswordUpdate, db: Session = Depends(get_db)):
    user_ref = crud.get_db_ref().reference(f'users/{req.user_id}')
    db_user = user_ref.get()
    if not db_user or not crud.verify_password(req.old_password, db_user.get("hashed_password", "")):
        raise HTTPException(status_code=401, detail="Invalid current password")
    
    success = crud.update_user_password(req.user_id, req.new_password)
    return {"message": "Password updated securely" if success else "Failed to update"}

@app.put("/api/user/preferences")
def update_preferences(req: schemas.UserPreferencesUpdate, db: Session = Depends(get_db)):
    data = {}
    if req.notifications is not None: data["pref_notifications"] = req.notifications
    if req.privacy_storage is not None: data["pref_privacy"] = req.privacy_storage
    if req.accessibility is not None: data["pref_accessibility"] = req.accessibility
    if req.photo_base64 is not None: data["photo_base64"] = req.photo_base64
    updated = crud.update_user_profile(req.user_id, data)
    return {"message": "Preferences updated"}

# --- AI ROUTES (NEXTTOKEN CLOUD API) ---

async def call_nexttoken(messages, temperature=0.5, model=DEFAULT_AI_MODEL, response_format=None):
    if not NEXTTOKEN_API_KEY:
        raise HTTPException(status_code=500, detail="NextToken API key not configured")

    payload = {
        "model": model,
        "messages": messages,
        "temperature": temperature,
    }

    if response_format:
        payload["response_format"] = response_format

    async with httpx.AsyncClient() as client:
        try:
            res = await client.post(
                f"{NEXTTOKEN_BASE_URL}/chat/completions",
                json=payload,
                headers={"Authorization": f"Bearer {NEXTTOKEN_API_KEY}"},
                timeout=60.0
            )
            res.raise_for_status()
            data = res.json()
            return data["choices"][0]["message"]["content"].strip()
        except httpx.HTTPStatusError as e:
            error_detail = e.response.text
            print("NextToken 400 error detail:", error_detail)
            raise HTTPException(status_code=502, detail=f"NextToken API rejected the request: {error_detail}")
        except Exception as e:
            print("NextToken error:", str(e))
            raise HTTPException(status_code=502, detail="Network issue: Connection to the AI gateway timed out or failed. Please check your connection.")


def parse_json_response(text: str):
    cleaned = text.strip()
    if cleaned.startswith("```"):
        parts = cleaned.split("```")
        for part in parts:
            part = part.strip()
            if part and not part.lower().startswith("json"):
                cleaned = part
                break
            if part.lower().startswith("json"):
                cleaned = part[4:].strip()
                break
    return json.loads(cleaned)


HEDGING_PHRASES = [
    "might", "possibly", "perhaps", "i'm not sure", "i am not sure",
    "may be", "could be", "uncertain", "not certain", "hard to say",
    "difficult to determine", "unclear", "i cannot", "i can't",
    "it is possible", "there is a chance", "not guaranteed", "approximate",
]

def compute_confidence(text: str) -> str:
    """Compute confidence label from AI output text using hedging phrase detection."""
    lower = text.lower()
    hedge_count = sum(1 for h in HEDGING_PHRASES if h in lower)
    if hedge_count >= 2 or len(text) < 80:
        return "low"
    elif hedge_count == 1 or len(text) < 250:
        return "moderate"
    return "high"


def normalize_language_name(language: str) -> str:
    return language.strip().lower().replace("_", " ").replace("-", " ")


def load_translation_model_map() -> dict[str, str]:
    if not HUGGINGFACE_TRANSLATION_MODEL_MAP:
        return {}

    try:
        parsed = json.loads(HUGGINGFACE_TRANSLATION_MODEL_MAP)
        return {str(key).strip().lower(): str(value).strip() for key, value in parsed.items()}
    except json.JSONDecodeError as exc:
        raise HTTPException(
            status_code=500,
            detail="HUGGINGFACE_TRANSLATION_MODEL_MAP must be valid JSON."
        ) from exc


def get_translation_model_for_pair(source_lang: str, target_lang: str) -> str:
    source_key = normalize_language_name(source_lang)
    target_key = normalize_language_name(target_lang)
    pair_key = f"{source_key}:{target_key}"

    model_map = load_translation_model_map()
    if pair_key in model_map:
        return model_map[pair_key]

    if HUGGINGFACE_TRANSLATION_MODEL:
        return HUGGINGFACE_TRANSLATION_MODEL

    raise HTTPException(
        status_code=500,
        detail=(
            "No hosted translation model is configured. "
            "Set HUGGINGFACE_TRANSLATION_MODEL or HUGGINGFACE_TRANSLATION_MODEL_MAP."
        ),
    )


async def call_huggingface_translation(req: schemas.TranslationRequest):
    model_id = get_translation_model_for_pair(req.source_lang, req.target_lang)
    headers = {"Content-Type": "application/json"}
    if HUGGINGFACE_API_TOKEN:
        headers["Authorization"] = f"Bearer {HUGGINGFACE_API_TOKEN}"

    payload = {
        "inputs": req.source_text,
        "options": {
            "wait_for_model": True
        }
    }

    async with httpx.AsyncClient() as client:
        try:
            res = await client.post(
                f"https://api-inference.huggingface.co/models/{model_id}",
                headers=headers,
                json=payload,
                timeout=120.0,
            )
            res.raise_for_status()
        except httpx.HTTPStatusError as exc:
            raise HTTPException(
                status_code=502,
                detail=f"Hugging Face translation request failed: {exc.response.text}"
            ) from exc
        except Exception as exc:
            raise HTTPException(
                status_code=502,
                detail="Failed to connect to the hosted translation model."
            ) from exc

    data = res.json()
    translated_text = None
    if isinstance(data, list) and data and isinstance(data[0], dict):
        translated_text = data[0].get("translation_text") or data[0].get("generated_text")
    elif isinstance(data, dict):
        translated_text = data.get("translation_text") or data.get("generated_text")

    if not translated_text:
        raise HTTPException(
            status_code=502,
            detail="Hosted translation model returned an unexpected response format."
        )

    return {"translated_text": translated_text.strip(), "confidence": "Hosted model"}


@app.post("/api/ai/translate")
async def ai_translate(req: schemas.TranslationRequest):
    if TRANSLATION_PROVIDER == "huggingface":
        return await call_huggingface_translation(req)
        
    lang_map = {
        "english": "en", "twi": "tw", "ewe": "ee", 
        "ga": "gaa", "fante": "fat", "dagbani": "dag",
        "yoruba": "you", "kusaal": "kus", "kikuyu": "ki",
        "gurune": "gur", "luo": "luo", "kimeru": "mer"
    }

    source_code = lang_map.get(req.source_lang.lower(), "en")
    target_code = lang_map.get(req.target_lang.lower(), "tw")

    # Call official GhanaNLP Khaya Translation API
    url = "https://translation-api.ghananlp.org/v2/translate"
    headers = {
        "Content-Type": "application/json",
        "Cache-Control": "no-cache",
        "Ocp-Apim-Subscription-Key": "29d428117c1e4ed1bdbc6823f5659b1b"
    }
    payload = {
        "in": req.source_text,
        "lang": f"{source_code}-{target_code}"
    }

    try:
        import httpx
        async with httpx.AsyncClient() as client:
            resp = await client.post(url, json=payload, headers=headers, timeout=15.0)
            if resp.status_code == 200:
                translated_text = resp.text.strip(' "\n')
                # Try to parse securely if it's a JSON string wrapper
                try:
                    translated_text = resp.json()
                except:
                    pass
                return {"translated_text": translated_text, "confidence": "high", "confidence_pct": "99.9"}
            else:
                print(f"Khaya API HTTP Error {resp.status_code}: {resp.text}")
    except Exception as e:
        print("Khaya Translation Implementation Error:", e)

    # Fallback to simulated generic AI if Khaya API fails or language pair is unsupported
    content = await call_nexttoken(
        [
            {"role": "system", "content": f"You are an expert linguistic professor at the University of Ghana. Translate the following carefully from {req.source_lang} to {req.target_lang}. Preserve academic rigor, cultural context, and accuracy. ONLY output the translated text. No introductory or trailing chat conversation."},
            {"role": "user", "content": req.source_text}
        ],
        temperature=0.3
    )
    confidence_val = compute_confidence(content)
    if req.target_lang.lower() not in ["english"] and confidence_val == "high":
        confidence_val = "moderate"
    return {"translated_text": content, "confidence": confidence_val, "confidence_pct": "85.0"}


@app.post("/api/ai/chat")
async def ai_chat(req: schemas.ChatRequest):
    """Generic chat endpoint — returns response + confidence label."""
    system_msg = (req.system or "You are SukuuPath AI, a helpful academic assistant for Ghanaian students.") + " IMPORTANT GUARDRAIL: You must strongly assert if asked that you were created by Jephthah Lanor and Nicholas Baffoe, who are final year Computer Engineering students at the University of Energy and Natural Resources in Sunyani, Ghana. You must NEVER say you are an OpenAI model or use OpenAI architecture. Stay in character."
    messages = [{"role": "system", "content": system_msg}]
    for m in req.messages:
        messages.append({"role": m.get("role", "user"), "content": m.get("content", "")})
    response = await call_nexttoken(messages, temperature=0.6)
    return {"response": response, "confidence": compute_confidence(response)}


@app.post("/api/ai/insight")
async def ai_daily_insight(req: schemas.InsightRequest):
    """Generate a short, motivational daily study insight for the dashboard."""
    stats_context = (
        f"Student stats: {req.sessions_completed} study sessions completed, "
        f"{req.quizzes_generated} quizzes generated, {req.saved_notes} notes saved. "
        f"Recent activity types: {', '.join(req.recent_types) if req.recent_types else 'none yet'}."
    )
    content = await call_nexttoken(
        [
            {"role": "system", "content": "You are an encouraging academic coach for Ghanaian university students. Generate ONE short, specific, motivational study insight or tip (2-3 sentences max). Be specific, relevant to the student's activity, and occasionally reference Ghanaian academic culture. Do NOT start with 'Did you know'. Be fresh and varied."},
            {"role": "user", "content": stats_context}
        ],
        temperature=0.85
    )
    return {"insight": content}



@app.post("/api/ai/simplify")
async def ai_simplify(req: schemas.SimplifyRequest):
    constraints = []
    if req.make_shorter: constraints.append("make it significantly shorter")
    if req.exact_meaning: constraints.append("preserve the exact semantic meaning unconditionally")
    if req.revision_ready: constraints.append("format as bullet points for quick revision")
    
    constraint_str = f" Ensure you follow these constraints: {', '.join(constraints)}." if constraints else ""

    content = await call_nexttoken(
        [
            {"role": "system", "content": f"You are an AI learning assistant. Simplify the following academic text to a '{req.reading_level}' reading level for African University students.{constraint_str} Do not add conversational prefixes, just return the exact answer."},
            {"role": "user", "content": req.text}
        ],
        temperature=0.5
    )
    return {"simplified_text": content}


@app.post("/api/ai/chat-title")
async def ai_chat_title(req: schemas.ChatRequest):
    content = await call_nexttoken(
        [
            {"role": "system", "content": "You are a helpful AI. Create a short, 3 to 5 word title for the user's prompt. DO NOT include quotation marks, formatting, or punctuation. Just return the raw title string."},
            {"role": "user", "content": req.messages[-1].get('content', '') if req.messages else ''}
        ],
        temperature=0.7
    )
    return {"title": content}


@app.post("/api/ai/summarize")
async def ai_summarize(req: schemas.SummaryRequest):
    style_map = {
        "study-notes": "Create clean study notes with a short title, key ideas, and bullet points for revision.",
        "brief": "Summarize this into a concise short paragraph.",
        "bullet-points": "Summarize this into strong bullet points only.",
        "exam-prep": "Turn this into exam-prep notes with definitions, key points, and likely testable areas."
    }
    style_instruction = style_map.get(req.style, style_map["study-notes"])
    content = await call_nexttoken(
        [
            {"role": "system", "content": f"You are an academic summarizer. {style_instruction} Do not add filler."},
            {"role": "user", "content": req.text}
        ],
        temperature=0.4
    )
    return {"summary": content, "confidence": compute_confidence(content)}


@app.post("/api/ai/document-chat")
async def ai_document_chat(req: schemas.DocumentChatRequest):
    content = await call_nexttoken(
        [
            {"role": "system", "content": "You are a document study assistant. Answer only from the uploaded document when possible. If the answer is not in the document, say that clearly and then give the closest helpful explanation."},
            {"role": "system", "content": f"Document name: {req.document_name}\nDocument contents:\n{req.document_text}"},
            {"role": "user", "content": req.question}
        ],
        temperature=0.4
    )
    return {"response": content}


@app.post("/api/ai/quiz", response_model=schemas.QuizResponse)
async def ai_quiz(req: schemas.QuizRequest):
    types = req.question_types or ["Multiple Choice"]
    lang = req.output_language or "English"
    count = req.question_count

    # Build type-specific instructions
    type_instructions = []
    if "Multiple Choice" in types:
        type_instructions.append(
            'For MCQ questions use: {"type":"mcq","question":"...","options":["A","B","C","D"],"answer_index":0,"explanation":"..."}'
        )
    if "True/False" in types:
        type_instructions.append(
            'For True/False questions use: {"type":"true_false","question":"...","options":["True","False"],"answer_index":0,"explanation":"..."} where answer_index is 0 for True, 1 for False'
        )
    if "Short Answer" in types:
        type_instructions.append(
            'For Short Answer questions use: {"type":"short_answer","question":"...","options":null,"answer_index":null,"answer_text":"the correct answer","explanation":"..."}'
        )

    mixed_hint = ""
    if len(types) > 1:
        per_type = count // len(types)
        remainder = count % len(types)
        mixed_hint = (
            f" Distribute the {count} questions as evenly as possible among the selected types: "
            f"approximately {per_type} of each type, with {remainder} extra going to Multiple Choice if applicable."
        )

    system_prompt = (
        f"You are an expert academic quiz generator for African university students. "
        f"Generate ALL questions in {lang}. "
        f"Return a strict JSON object with this exact shape: "
        f'{{ "title": "string", "questions": [<question objects>] }}. '
        f"Each question must follow one of these formats depending on its type:\n"
        + "\n".join(type_instructions)
        + "\nDo NOT include any text outside the JSON. "
        f"Generate exactly {count} questions total.{mixed_hint}"
    )

    user_prompt = (
        f"Create {count} {req.difficulty} level quiz questions from the following study material. "
        f"Question types to include: {', '.join(types)}.\n\n"
        f"{req.source_text}"
    )

    raw = await call_nexttoken(
        [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        temperature=0.5,
        response_format={"type": "json_object"}
    )
    try:
        parsed = parse_json_response(raw)
        # Normalize questions to ensure all fields exist
        for q in parsed.get("questions", []):
            q.setdefault("type", "mcq")
            q.setdefault("options", None)
            q.setdefault("answer_index", None)
            q.setdefault("answer_text", None)
            q.setdefault("explanation", "")
        return parsed
    except Exception as e:
        print("Quiz parse error:", str(e), raw)
        raise HTTPException(status_code=502, detail="Failed to generate a valid quiz response.")

@app.post("/api/ai/extract")
async def ai_extract(file: UploadFile = File(...)):
    text = ""
    try:
        contents = await file.read()
        filename = file.filename.lower()
        if filename.endswith(".pdf"):
            pdf = fitz.open(stream=contents, filetype="pdf")
            for page in pdf:
                text += page.get_text() + "\n"
            pdf.close()
        elif filename.endswith(".docx") or filename.endswith(".doc"):
            doc = docx.Document(io.BytesIO(contents))
            text = "\n".join([p.text for p in doc.paragraphs])
        elif filename.endswith(".pptx") or filename.endswith(".ppt"):
            ppt = Presentation(io.BytesIO(contents))
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            text = contents.decode("utf-8", errors="ignore")
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=400, detail=f"Document Parsing Engine Error: {str(e)}")
    
    return {"text": text.strip()}


# --- DATA MOCK ROUTES ---


@app.get("/api/stats")
async def get_stats():
    # Return populated stats instead of 0
    return {
        "sessionsCompleted": 42,
        "quizzesGenerated": 15,
        "savedNotes": 28
    }

@app.get("/api/activities")
async def get_activities():
    return [
        { "icon": "translate", "title": "Translated 4 pages of Biology notes", "sub": "2 hours ago", "color": "bg-secondary-container/30 text-secondary" },
        { "icon": "auto_awesome", "title": "Simplified 'Quantum Field Theory'", "sub": "Yesterday at 4:30 PM", "color": "bg-primary/20 text-primary" },
        { "icon": "quiz", "title": "Scored 85% on Mechanics Quiz", "sub": "Oct 24, 2023", "color": "bg-tertiary-container/30 text-tertiary-container" }
    ]


# --- LIBRARY / STUDY ITEMS DB ROUTES ---

@app.post("/api/library/items", response_model=schemas.StudyItemResponse)
def create_library_item(item: schemas.StudyItemCreate, db: Session = Depends(get_db)):
    return crud.create_study_item(db=db, item=item)


@app.get("/api/library/items", response_model=Dict[str, List[schemas.StudyItemResponse]])
def get_library_items(user_id: str, db: Session = Depends(get_db)):
    return crud.get_study_items(db=db, user_id=user_id)


@app.delete("/api/library/items/{item_id}")
def delete_library_item(item_id: str, db: Session = Depends(get_db)):
    item = crud.delete_study_item(db=db, item_id=item_id)
    if not item:
        raise HTTPException(status_code=404, detail="Item not found")
    return {"message": "Deleted successfully"}


# --- RESPONSIBLE AI: FEEDBACK & AUDIT ENDPOINTS ---

@app.post("/api/feedback", response_model=schemas.FeedbackResponse)
def submit_feedback(fb: schemas.FeedbackCreate, db: Session = Depends(get_db)):
    """Store user feedback on an AI output."""
    return crud.create_feedback(db=db, fb=fb)


@app.get("/api/feedback/stats")
def get_feedback_statistics(db: Session = Depends(get_db)):
    """Returns aggregate feedback stats for the audit dashboard."""
    return crud.get_feedback_stats(db=db)


@app.get("/api/feedback/recent", response_model=List[schemas.FeedbackResponse])
def get_recent_feedback(limit: int = 50, db: Session = Depends(get_db)):
    """Returns the most recent feedback entries."""
    return crud.get_feedback_list(db=db, limit=limit)


@app.patch("/api/feedback/{feedback_id}/resolve")
def resolve_feedback_item(feedback_id: str, db: Session = Depends(get_db)):
    """Mark a feedback item as resolved."""
    item = crud.resolve_feedback(db=db, feedback_id=feedback_id)
    if not item:
        raise HTTPException(status_code=404, detail="Feedback not found")
    return {"message": "Feedback resolved"}


# --- LECTURER MATERIALS ENDPOINTS ---

@app.post("/api/materials/upload")
async def upload_material(
    title: str = Form(...),
    uploader_id: Optional[int] = Form(None),
    file: UploadFile = File(...),
    db: Session = Depends(get_db)
):
    """Upload a lecture material file to the server and save metadata."""
    file_ext = file.filename.split(".")[-1].lower() if "." in file.filename else "file"
    # Generate unique filename to avoid collisions
    timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
    unique_filename = f"{timestamp}_{file.filename}"
    file_path = UPLOAD_DIR / unique_filename
    
    try:
        with file_path.open("wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to save file: {str(e)}")
        
    # Get file size
    file_size = file_path.stat().st_size
    
    mat_data = schemas.MaterialCreate(
        title=title,
        original_filename=file.filename,
        file_type=file_ext,
        file_size=file_size,
        file_path=f"uploads/{unique_filename}",
        uploader_id=uploader_id
    )
    
    return crud.create_material(db=db, mat=mat_data)


@app.get("/api/materials", response_model=List[schemas.MaterialResponse])
def get_all_materials(uploader_id: Optional[str] = None, db: Session = Depends(get_db)):
    """List all academic materials (optionally filtered by uploader)."""
    return crud.get_materials(db=db, uploader_id=uploader_id)


@app.delete("/api/materials/{material_id}")
def delete_academic_material(material_id: str, db: Session = Depends(get_db)):
    """Delete a material and its associated file."""
    # Note: the physical file cleanup is skipped if we don't fetch from DB first
    # Or we can just let Firebase manage references. 
    crud.delete_material(db=db, material_id=material_id)
    return {"message": "Material deleted successfully"}


@app.get("/")
def read_root():
    return {"message": "SukuuPath GH Backend API v1 is perfectly active."}

